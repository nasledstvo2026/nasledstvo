#!/usr/bin/env python3
"""Генератор презентации: Отдел сопровождения дирекции радиоэлектроники"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Цветовая схема ──────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x0B, 0x1F, 0x3F)   # тёмно-синий фон
BLUE_MID    = RGBColor(0x1A, 0x3A, 0x6C)   # акцентный синий
BLUE_ACCENT = RGBColor(0x2D, 0x5F, 0x9E)   # светлый акцент
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT  = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_TEXT   = RGBColor(0x99, 0x99, 0x99)
GREEN_KPI   = RGBColor(0x2E, 0xCC, 0x71)
YELLOW_KPI  = RGBColor(0xF3, 0x9C, 0x12)
ORANGE_KPI  = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ── Хелперы ─────────────────────────────────────────────────────
def add_bg(slide, color=BLUE_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_list(tf, items, font_size=14, color=WHITE, spacing=Pt(8)):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0

def add_decor_line(slide, left, top, width, color=BLUE_ACCENT, height=Pt(3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_kpi_tile(slide, left, top, width, height, value, label, color=GREEN_KPI):
    """Плитка KPI: крупная цифра + подпись."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x11, 0x2B, 0x50)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_LIGHT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

def add_footer(slide, text="Отдел сопровождения дирекции радиоэлектроники | Конфиденциально"):
    add_textbox(slide, 0.5, 7.0, 8, 0.4, text, font_size=9, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 1 — Титульный
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BLUE_DARK)

add_textbox(slide, 1.0, 1.5, 11, 1.2,
            "ОТДЕЛ СОПРОВОЖДЕНИЯ", font_size=44, bold=True, color=WHITE)
add_textbox(slide, 1.0, 2.6, 11, 1.0,
            "ДИРЕКЦИИ РАДИОЭЛЕКТРОНИКИ", font_size=36, bold=False, color=BLUE_ACCENT)

add_decor_line(slide, 1.0, 3.7, 3.0, BLUE_ACCENT)

add_textbox(slide, 1.0, 4.1, 8, 0.6,
            "Операционная поддержка бизнес-процессов", font_size=20, color=GRAY_LIGHT)
add_textbox(slide, 1.0, 6.5, 8, 0.4,
            "2026 | Для внутреннего пользования", font_size=12, color=GRAY_TEXT)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 2 — Миссия и цель
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 10, 0.8, "МИССИЯ ОТДЕЛА", font_size=30, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 2.5, BLUE_ACCENT)

tf = add_textbox(slide, 0.8, 1.6, 5.5, 4.5, "", font_size=16, color=WHITE)
add_bullet_list(tf, [
    "🎯 Централизованное сервисное звено дирекции радиоэлектроники",
    "⚙️ Бесперебойное обеспечение административных, закупочных и финансовых процессов",
    "🔗 Единая точка входа для смежных управлений",
    "📋 Стандартизация процедур — снижение операционных рисков",
    "⏱️ Оперативное закрытие потребностей без задержек",
], font_size=15, color=WHITE)

# Три плитки ценностей
values = [
    ("СКОРОСТЬ", "Минимальный срок\nобработки заявок и\nоформления документов"),
    ("НАДЁЖНОСТЬ", "Гарантированное\nисполнение каждого\nэтапа процессов"),
    ("КОНТРОЛЬ", "Полный аудит ЖЦ\nкаждой закупки и\nкомандировки"),
]
for idx, (title, desc) in enumerate(values):
    x = 7.2 + idx * 2.0
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(1.8), Inches(2.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x11, 0x2B, 0x50); shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = BLUE_ACCENT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(10); p2.font.color.rgb = GRAY_LIGHT
    p2.alignment = PP_ALIGN.CENTER

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 3 — Закупочная деятельность
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "УПРАВЛЕНИЕ ЗАКУПКАМИ И ПОСТАВЩИКАМИ", font_size=28, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 3.5, BLUE_ACCENT)

# Схема: три шага
steps = [
    ("01", "ЗАЯВКА", "Размещение заявок\nна закупку\nв системе"),
    ("02", "НОМЕНКЛАТУРА", "Заведение и ведение\nноменклатуры\nдля закупок"),
    ("03", "ДОГОВОР", "Заключение договоров\nс поставщиками\nи контрагентами"),
]
for idx, (num, title, desc) in enumerate(steps):
    x = 1.5 + idx * 3.8
    # Номер
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(1.8), Inches(1.0), Inches(1.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_ACCENT; shape.line.fill.background()
    tf = shape.text_frame; p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    # Заголовок
    add_textbox(slide, x, 3.1, 2.8, 0.5, title, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Описание
    add_textbox(slide, x, 3.6, 2.8, 1.2, desc, font_size=12, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

# Нижний блок
tf = add_textbox(slide, 0.8, 5.2, 12, 1.5, "", font_size=14, color=WHITE)
add_bullet_list(tf, [
    "📊 Единый реестр контрагентов и договорной базы",
    "🔄 Полный цикл: от потребности до исполнения договора",
    "📋 Стандартизированные шаблоны договорной документации",
], font_size=14, color=WHITE)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 4 — Командировочное обеспечение
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "ОРГАНИЗАЦИЯ КОМАНДИРОВОК", font_size=28, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 3.0, BLUE_ACCENT)

# Таймлайн
timeline = [
    ("📝", "Заявка в 1С", "Оформление и запуск\nкомандировки в системе"),
    ("🎫", "Выкуп билетов\nи гостиниц", "Бронирование транспорта\nи проживания"),
    ("✅", "Закрытие", "Полный пакет отчётных\nдокументов"),
]
for idx, (icon, title, desc) in enumerate(timeline):
    y = 2.0 + idx * 1.8
    # Линия слева
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.6), Inches(y + 0.3), Inches(0.03), Inches(1.4))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_ACCENT; shape.line.fill.background()
    # Кружок
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), Inches(y + 0.15), Inches(0.6), Inches(0.6))
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_ACCENT; shape.line.fill.background()
    tf = shape.text_frame; p = tf.paragraphs[0]; p.text = icon; p.font.size = Pt(18); p.alignment = PP_ALIGN.CENTER
    # Текст
    add_textbox(slide, 2.3, y, 2.5, 0.5, title, font_size=16, bold=True, color=WHITE)
    add_textbox(slide, 2.3, y + 0.5, 4.5, 1.0, desc, font_size=12, color=GRAY_LIGHT)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 5 — Авансовые отчёты
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "ФИНАНСОВО-ОТЧЁТНАЯ ДИСЦИПЛИНА", font_size=28, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 3.2, BLUE_ACCENT)

# Две колонки
cols = [
    ("💰", "СОБСТВЕННЫЕ СРЕДСТВА", [
        "Формирование авансовых отчётов сотрудников",
        "Контроль первичной документации",
        "Соблюдение сроков закрытия",
        "Утверждение отчётов в системе",
    ]),
    ("🏭", "ЗАКУПКА ПРОИЗВОДСТВА", [
        "Закрытие авансовых отчётов для нужд производства",
        "Привязка к ЖЦ закупки",
        "Сверка с договорной базой",
        "Финансовый контроль расхода средств",
    ]),
]
for idx, (icon, title, items) in enumerate(cols):
    x = 1.2 + idx * 5.8
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(5.2), Inches(4.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x11, 0x2B, 0x50); shape.line.fill.background()
    add_textbox(slide, x + 0.4, 1.8, 4.5, 0.5, f"{icon}  {title}", font_size=16, bold=True, color=BLUE_ACCENT)
    tf = add_textbox(slide, x + 0.4, 2.5, 4.3, 3.0, "", font_size=13, color=WHITE)
    add_bullet_list(tf, items, font_size=13, color=WHITE)

# Нижняя строка
add_textbox(slide, 0.8, 6.7, 12, 0.3,
            "🔄 Полный цикл: от выдачи аванса до утверждения отчёта",
            font_size=12, color=GRAY_LIGHT)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 6 — Административное обеспечение
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "АДМИНИСТРАТИВНО-ХОЗЯЙСТВЕННАЯ ПОДДЕРЖКА", font_size=28, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 3.8, BLUE_ACCENT)

tiles = [
    ("🪪", "ПРОПУСКНАЯ СИСТЕМА", "Оформление и учёт пропусков\nКонтроль сроков действия\nЕдиная база выданных пропусков"),
    ("🛒", "АДМИНИСТРАТИВНЫЕ ЗАКУПКИ", "Закупка хозяйственных нужд\nФормирование потребностей для УР\nПланирование и бюджетирование"),
    ("💻", "ОФИСНАЯ ТЕХНИКА", "Оснащение новых сотрудников\nПолный комплект в день выхода\nУчёт и инвентаризация оборудования"),
]
for idx, (icon, title, desc) in enumerate(tiles):
    x = 1.2 + idx * 3.9
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(3.5), Inches(4.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x11, 0x2B, 0x50); shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{icon}  {title}"; p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = GRAY_LIGHT
    p2.alignment = PP_ALIGN.CENTER

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 7 — Взаимодействие
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "КОНТУР ВЗАИМОДЕЙСТВИЯ", font_size=28, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 3.0, BLUE_ACCENT)

# Две колонки
cols_data = [
    ("🏢 ВНУТРЕННИЕ СТЕЙКХОЛДЕРЫ", [
        "Смежные управления",
        "Дирекция радиоэлектроники",
        "Подразделения УР",
        "Бухгалтерия и финансовый отдел",
    ]),
    ("🌐 ВНЕШНИЕ КОНТРАГЕНТЫ", [
        "Поставщики оборудования и услуг",
        "Транспортные компании",
        "Гостиницы и операторы бронирования",
        "Контролирующие органы",
    ]),
]
for idx, (title, items) in enumerate(cols_data):
    x = 1.2 + idx * 5.8
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(5.2), Inches(3.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x11, 0x2B, 0x50); shape.line.fill.background()
    add_textbox(slide, x + 0.3, 1.8, 4.5, 0.5, title, font_size=16, bold=True, color=BLUE_ACCENT)
    tf = add_textbox(slide, x + 0.3, 2.5, 4.5, 2.0, "", font_size=13, color=WHITE)
    add_bullet_list(tf, items, font_size=13, color=WHITE)

# Коммуникационный блок снизу
add_textbox(slide, 0.8, 5.5, 12, 0.4,
            "🔄 Регламентированные цепочки согласования   |   📞 Единое окно для внешних запросов",
            font_size=14, bold=True, color=WHITE)

add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 8 — KPI и метрики (ДОКРУЧЕННЫЙ)
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "КЛЮЧЕВЫЕ ПОКАЗАТЕЛИ ЭФФЕКТИВНОСТИ", font_size=28, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 3.8, BLUE_ACCENT)

# Ряд 1 — операционные KPI (количественные)
kpi_row1 = [
    ("< 3 дней", "Средний срок обработки\nзаявки на закупку", GREEN_KPI),
    ("≥ 95%", "Доля командировок,\nзакрытых в срок", GREEN_KPI),
    ("> 1 200 ед.", "Объём документооборота\nв месяц", BLUE_ACCENT),
    ("80+", "Активных договоров\nс поставщиками", BLUE_ACCENT),
]
for idx, (value, label, color) in enumerate(kpi_row1):
    x = 0.8 + idx * 3.2
    add_kpi_tile(slide, x, 1.6, 2.8, 2.0, value, label, color)

# Ряд 2 — SLA и качественные KPI
kpi_row2 = [
    ("< 4 часов", "SLA оформления\nпропуска", GREEN_KPI),
    ("100%", "Обеспеченность техникой\nновых сотрудников в день выхода", GREEN_KPI),
    ("< 5 дней", "Средний срок\nзакрытия авансового отчёта", YELLOW_KPI),
    ("0", "Просроченных\nдоговоров поставки", GREEN_KPI),
]
for idx, (value, label, color) in enumerate(kpi_row2):
    x = 0.8 + idx * 3.2
    add_kpi_tile(slide, x, 4.0, 2.8, 2.0, value, label, color)

# Подвал с пояснением
add_textbox(slide, 0.8, 6.5, 12, 0.4,
            "📊 Целевые значения KPI обновляются ежеквартально на основе ретроспективного анализа",
            font_size=11, color=GRAY_TEXT)
add_footer(slide)

# ══════════════════════════════════════════════════════════════════
# СЛАЙД 9 — Заключение и контакты
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_textbox(slide, 0.8, 0.4, 11, 0.8, "КОНТАКТЫ", font_size=30, bold=True, color=WHITE)
add_decor_line(slide, 0.8, 1.1, 2.0, BLUE_ACCENT)

add_textbox(slide, 0.8, 2.0, 11, 0.8,
            "Отдел сопровождения — ваш операционный фундамент",
            font_size=24, bold=False, color=WHITE)

contacts = [
    "📧 Электронная почта",
    "📞 Телефон",
    "📍 Расположение",
]
for idx, line in enumerate(contacts):
    add_textbox(slide, 0.8, 3.2 + idx * 0.6, 6, 0.5, line, font_size=16, color=GRAY_LIGHT)

# CTA-блок
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.5), Inches(5.0), Inches(2.0))
shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_ACCENT; shape.line.fill.background()
tf = shape.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "ГОТОВЫ ОБЕСПЕЧИТЬ"; p.font.size = Pt(22); p.font.bold = True
p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "ВАШ ЗАПРОС В ДЕНЬ ОБРАЩЕНИЯ"; p2.font.size = Pt(16)
p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER

add_footer(slide)

# ── Сохранение ───────────────────────────────────────────────────
output_path = os.path.expanduser("~/.openclaw/workspace/Отдел_сопровождения_РЭ.pptx")
prs.save(output_path)
print(f"✅ Презентация сохранена: {output_path}")
