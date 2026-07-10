#!/usr/bin/env python3
"""Генерация pptx-презентации «Зарплатный монитор региона» для руководителя отдела продаж Сбера."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "zarplatny-monitor-sber.pptx")
OUTPUT = os.path.abspath(OUTPUT)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK = RGBColor(0x0A, 0x0E, 0x14)
SBER_GREEN = RGBColor(0x21, 0xA0, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
LIGHT_GRAY = RGBColor(0x8B, 0x94, 0x9E)
CARD_BG = RGBColor(0x16, 0x1B, 0x22)
BORDER = RGBColor(0x21, 0x26, 0x2D)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
RED = RGBColor(0xFF, 0x45, 0x45)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=CARD_BG, line_color=BORDER):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=WHITE, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(8)
    return txBox

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_shape(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(4.5))
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1), "Зарплатный монитор региона", font_size=44, bold=True, color=SBER_GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1), "Цифровой инструмент разведки для регионального руководителя\nпродаж «Зарплатного проекта» Сбера", font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6), "Июль 2026", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 2: Проблема ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8), fill_color=CARD_BG)
add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7), "Проблема", font_size=32, bold=True, color=SBER_GREEN)

cards = [
    (0.5, "📋", "Холодные списки\nНет актуальной базы новых юрлиц — обзвон вслепую. Базы устаревают за недели."),
    (4.3, "🔍", "Слепые переговоры\nНе знаешь, у кого клиент сейчас. ВТБ? Т-Банк? Альфа? Аргументов нет."),
    (8.1, "⏰", "Упущенные тендеры\nГосзакупки на банковские услуги проходят без вашего участия — узнаёте постфактум."),
]

for left, emoji, text in cards:
    add_shape(slide, Inches(left), Inches(1.8), Inches(3.5), Inches(3.5))
    add_textbox(slide, Inches(left + 0.3), Inches(2.0), Inches(3), Inches(0.6), emoji, font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(left + 0.3), Inches(2.8), Inches(3), Inches(2.3), text, font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 3: Решение ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8), fill_color=CARD_BG)
add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7), "Решение: «Зарплатный монитор»", font_size=32, bold=True, color=SBER_GREEN)

modules = [
    (0.5, "📦", "Новые юрлица", "Ежедневный поток свежих ООО и ИП с контактами. Готовая база для обзвона sales-команды."),
    (3.5, "🏪", "Свободные предприятия", "Тендеры на банковские услуги — кто ищет. Пока контракт не подписан, можно зайти."),
    (6.5, "📰", "Конкурентная разведка", "Негатив о конкурентах (задержки, сбои, рост тарифов) — готовые аргументы для переговоров."),
    (9.5, "📊", "Статистика региона", "Динамика регистраций, проникновение по районам, топ-крупных клиентов."),
]

for left, emoji, title, desc in modules:
    add_shape(slide, Inches(left), Inches(1.8), Inches(2.7), Inches(4.5))
    add_textbox(slide, Inches(left + 0.2), Inches(2.0), Inches(2.3), Inches(0.5), emoji, font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(left + 0.2), Inches(2.6), Inches(2.3), Inches(0.5), title, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(left + 0.2), Inches(3.2), Inches(2.3), Inches(2.5), desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 4: Как это работает ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8), fill_color=CARD_BG)
add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7), "Как это работает", font_size=32, bold=True, color=SBER_GREEN)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.5), "Автоматический пайплайн — без участия человека", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Pipeline boxes
pipeline = [
    (0.5, "ФНС России", "Новые юрлица\nрегиона (ОКТМО)"),
    (3.0, "ЕИС\n(zakupki.gov.ru)", "Тендеры на\nбанк. услуги"),
    (5.5, "Поиск СМИ", "Упоминания\nбанков в регионе"),
    (8.0, "⚙️ Сборщик", "JSON → HTML +\nTelegram-сводка"),
    (10.5, "👤 Руководителю", "Telegram 08:00\n+ сайт компании"),
]

for left, title, desc in pipeline:
    add_shape(slide, Inches(left), Inches(2.8), Inches(2.2), Inches(2.5))
    add_textbox(slide, Inches(left + 0.1), Inches(3.0), Inches(2.0), Inches(0.8), title, font_size=13, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(left + 0.1), Inches(3.8), Inches(2.0), Inches(1.2), desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5), "Расписание: ежедневно до 08:00 (будни)", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 5: Пример сводки ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8), fill_color=CARD_BG)
add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7), "Пример сводки в Telegram", font_size=32, bold=True, color=SBER_GREEN)

add_shape(slide, Inches(2.5), Inches(1.8), Inches(8.3), Inches(4.5))
msg_lines = [
    "☀️ Зарплатный монитор — Москва, 01.07.2026",
    "",
    "📦 Новые юрлица: 14 (из них 9 — целевые, 5+ сотрудников)",
    "   ООО «ТехноСервис» — ИНН 7712345678 — оптовая торговля",
    "   ООО «СтройГрупп» — ИНН 7743123456 — строительство",
    "   Контакты: +7 (495) 123-45-67, techno@mail.ru",
    "",
    "🏦 Тендеры ЗП: 3 активных, 2 завершены",
    "   МБОУ СОШ №45 — НМЦК 1,2 млн ₽ — приём до 15.07",
    "   Победитель: ВТБ (можно заходить с КП)",
    "",
    "📰 Негатив о конкурентах:",
    "   • ВТБ — задержки выпуска зарплатных карт (banki.ru)",
    "   • Т-Банк — сбой в день зарплаты (cnews.ru)",
    "",
    "📍 Карта: 5 из 10 районов охвачены",
]
add_bullet_textbox(slide, Inches(3), Inches(2.0), Inches(7.5), Inches(4), msg_lines, font_size=13, color=LIGHT_GRAY)

# ── Slide 6: Данные ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8), fill_color=CARD_BG)
add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7), "Источники данных — всё из открытых источников", font_size=32, bold=True, color=SBER_GREEN)

sources = [
    (0.5, "ФНС России (egrul.nalog.ru)", "ЕГРЮЛ/ЕГРИП — данные о регистрации,\nадресах, ОКВЭД, контактах юрлиц\n✅ Бесплатно, без ключей"),
    (4.3, "ЕИС (zakupki.gov.ru)", "Тендеры по 44-ФЗ / 223-ФЗ — банковские\nуслуги, РКО, зарплатные проекты\n✅ Открытые данные"),
    (8.1, "SearXNG / web_search", "Мониторинг СМИ — упоминания банков,\nтарифов, отзывов, сбоев\n✅ Уже настроено"),
]

for left, title, desc in sources:
    add_shape(slide, Inches(left), Inches(1.8), Inches(3.5), Inches(3.5))
    add_textbox(slide, Inches(left + 0.2), Inches(2.0), Inches(3.1), Inches(0.6), title, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(left + 0.2), Inches(2.7), Inches(3.1), Inches(2.3), desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5), "Никаких платных API. Только официальные открытые данные и собственные парсеры.", font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ── Slide 7: Стоимость ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8), fill_color=CARD_BG)
add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7), "Стоимость и сроки", font_size=32, bold=True, color=SBER_GREEN)

add_shape(slide, Inches(0.5), Inches(1.8), Inches(3.8), Inches(4.5))
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(3.2), Inches(0.5), "Срок внедрения", font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(2.7), Inches(3.2), Inches(3), "MVP: 2-3 дня\n\n1 день — парсер ФНС\n1 день — адаптация ЕИС\n0,5 дня — сборка отчёта\n0,5 дня — настройка cron + Telegram", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.8), Inches(1.8), Inches(3.8), Inches(4.5))
add_textbox(slide, Inches(5.1), Inches(2.0), Inches(3.2), Inches(0.5), "Себестоимость", font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(5.1), Inches(2.7), Inches(3.2), Inches(3), "0 ₽ в месяц\n\nТолько токены LLM:\n~200 запросов/мес\n~$50-100\nна DeepSeek API", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(9.1), Inches(1.8), Inches(3.8), Inches(4.5))
add_textbox(slide, Inches(9.4), Inches(2.0), Inches(3.2), Inches(0.5), "Выгода", font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.4), Inches(2.7), Inches(3.2), Inches(3), "Бесплатный источник\nготовых лидов каждый день\n\nЭкономия времени\nменеджеров на поиск\n~20 часов/неделю", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 8: Демо ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.8), fill_color=CARD_BG)
add_textbox(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7), "Что мы можем показать уже завтра", font_size=32, bold=True, color=SBER_GREEN)

add_shape(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.5))
add_textbox(slide, Inches(2), Inches(2.0), Inches(9), Inches(1), "🗺️ Выберите регион. Завтра в 08:00 — первый Telegram-дайджест.\nЧерез 3 дня — полноценный HTML-отчёт на сайте.\nНикаких обязательств — смотрим, работает ли в вашем регионе.", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(2.5), Inches(4.0), Inches(8.3), Inches(2.5))
add_textbox(slide, Inches(3), Inches(4.2), Inches(7.5), Inches(0.5), "Готовы попробовать?", font_size=24, bold=True, color=SBER_GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3), Inches(5.0), Inches(7.5), Inches(1.2), "Достаточно назвать регион — и мы запускаем.\n\nКонтакт: @Kirill_syst", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Slide 9: Контакты ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_shape(slide, Inches(2), Inches(2), Inches(9.3), Inches(3.5))
add_textbox(slide, Inches(2.5), Inches(2.3), Inches(8), Inches(0.8), "Спасибо за внимание", font_size=36, bold=True, color=SBER_GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.5), Inches(3.3), Inches(8), Inches(0.5), "Зарплатный монитор региона — разведданные для роста продаж", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(2.5), Inches(4.2), Inches(8), Inches(0.5), "Telegram: @Kirill_syst", font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ── Save ──
prs.save(OUTPUT)
print(f"✅ Презентация сохранена: {OUTPUT}")
print(f"   Размер: {os.path.getsize(OUTPUT) / 1024:.0f} KB")
